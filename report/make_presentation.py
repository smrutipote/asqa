"""Generate ASQA presentation PPTX."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
FIG = HERE / "figures"
OUT = HERE / "ASQA_presentation.pptx"

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x0E, 0x7C, 0x86)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
GREY = RGBColor(0x55, 0x5B, 0x66)
ACCENT = RGBColor(0xE2, 0x8D, 0x1F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=NAVY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"\u2022  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"


def header(slide, title, subtitle=None):
    add_rect(slide, Emu(0), Emu(0), SW, Inches(1.1), NAVY)
    add_rect(slide, Emu(0), Inches(1.1), SW, Inches(0.08), TEAL)
    add_text(slide, Inches(0.5), Inches(0.22), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.72), Inches(12), Inches(0.4),
                 subtitle, size=14, color=RGBColor(0xCF, 0xE3, 0xE7))
    add_text(slide, Inches(12.3), Inches(7.15), Inches(0.9), Inches(0.25),
             "ASQA", size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ---------- 1. Title ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, Emu(0), Emu(0), SW, SH, NAVY)
add_rect(s, Inches(0), Inches(3.0), SW, Inches(0.05), TEAL)
add_rect(s, Inches(0), Inches(3.4), SW, Inches(0.02), ACCENT)
add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.0),
         "ASQA", size=72, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.8),
         "Autonomous Software Quality Assurance", size=28, color=RGBColor(0xCF, 0xE3, 0xE7))
add_text(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.6),
         "A Multi-Agent LLM Pipeline for Bug Detection,", size=22, color=WHITE)
add_text(s, Inches(0.7), Inches(4.1), Inches(12), Inches(0.6),
         "Test Generation, and Fix Suggestion", size=22, color=WHITE)
add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.4),
         "Smruti Pote   |   x24269522", size=18, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.4),
         "MSc Artificial Intelligence  \u2014  National College of Ireland", size=14,
         color=RGBColor(0xCF, 0xE3, 0xE7))

# ---------- 2. Agenda ----------
s = prs.slides.add_slide(BLANK)
header(s, "Agenda")
items = [
    ("1", "Problem & motivation"),
    ("2", "Research question"),
    ("3", "ASQA architecture & agents"),
    ("4", "Evaluation setup"),
    ("5", "Results vs. baselines"),
    ("6", "Implications"),
    ("7", "Limitations & future work"),
]
col_w = Inches(5.8)
row_h = Inches(0.65)
for i, (n, t) in enumerate(items):
    col = i // 4
    row = i % 4
    x = Inches(0.7) + col * Inches(6.0)
    y = Inches(1.6) + row * Inches(0.9)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    tf = circle.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = n
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
    add_text(s, x + Inches(0.85), y + Inches(0.1), col_w, Inches(0.5),
             t, size=20, color=NAVY)

# ---------- 3. Problem / Motivation ----------
s = prs.slides.add_slide(BLANK)
header(s, "The Problem", "Why automate software quality assurance?")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(7.5), Inches(5), [
    "Testing and debugging consume 25\u201350% of project budgets.",
    "Poor software quality cost the US $2.41T in 2022 (CISQ).",
    "Single-agent LLM approaches resolve only 3.97% of SWE-bench issues.",
    "Context-window, recovery, and multi-task limits drive failures.",
    "Opportunity: decompose QA into specialised, coordinated agents.",
], size=20)

# Callout card
card_x, card_y, card_w, card_h = Inches(8.6), Inches(1.8), Inches(4.2), Inches(4.5)
add_rect(s, card_x, card_y, card_w, card_h, LIGHT)
add_rect(s, card_x, card_y, Inches(0.12), card_h, TEAL)
add_text(s, card_x + Inches(0.3), card_y + Inches(0.25), card_w, Inches(0.5),
         "Context", size=14, bold=True, color=TEAL)
add_text(s, card_x + Inches(0.3), card_y + Inches(0.75), card_w - Inches(0.4), Inches(0.6),
         "$2.41T", size=40, bold=True, color=NAVY)
add_text(s, card_x + Inches(0.3), card_y + Inches(1.55), card_w - Inches(0.4), Inches(0.6),
         "Annual US cost of poor\nsoftware quality (CISQ, 2022)", size=14, color=GREY)
add_text(s, card_x + Inches(0.3), card_y + Inches(2.6), card_w - Inches(0.4), Inches(0.6),
         "3.97%", size=40, bold=True, color=ACCENT)
add_text(s, card_x + Inches(0.3), card_y + Inches(3.4), card_w - Inches(0.4), Inches(0.6),
         "Single-agent SWE-bench\nresolution rate", size=14, color=GREY)

# ---------- 4. Research Question ----------
s = prs.slides.add_slide(BLANK)
header(s, "Research Question")
box_x, box_y, box_w, box_h = Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.6)
add_rect(s, box_x, box_y, box_w, box_h, LIGHT)
add_rect(s, box_x, box_y, box_w, Inches(0.12), TEAL)
add_text(s, box_x + Inches(0.6), box_y + Inches(0.4), Inches(10.5), Inches(2.2),
         ("Will a multi-agent LLM pipeline, specialised in QA subtasks and "
          "coordinated by LangGraph, achieve a Bug Detection Rate > 60% and "
          "a Mean Time to Report < 120 seconds \u2014 outperforming single-agent "
          "baselines?"),
         size=22, bold=True, color=NAVY)

# Target chips
targets = [("BDR > 60%", TEAL), ("MTTR < 120 s", ACCENT), ("Reliability", NAVY)]
for i, (t, c) in enumerate(targets):
    x = Inches(1.3) + i * Inches(3.9)
    y = Inches(5.3)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.5), Inches(1.1))
    chip.fill.solid(); chip.fill.fore_color.rgb = c
    chip.line.fill.background()
    tf = chip.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE

# ---------- 5. Related Work ----------
s = prs.slides.add_slide(BLANK)
header(s, "Related Work \u2014 Gap", "Where ASQA fits")
rows = [
    ("Codex, GPT-4, Claude 3.5", "Strong code generation, still ~19% buggy output"),
    ("Chain-of-Thought, ReAct", "Decomposition improves reasoning"),
    ("AgentCoder (3 agents)", "77.4% pass@1 on HumanEval \u2014 code gen only"),
    ("Defects4J / BugsInPy / SWE-bench", "Standard benchmarks for faults & issues"),
]
y = Inches(1.6)
for i, (k, v) in enumerate(rows):
    row_y = y + i * Inches(0.85)
    add_rect(s, Inches(0.7), row_y, Inches(4.3), Inches(0.75),
             LIGHT if i % 2 == 0 else WHITE)
    add_rect(s, Inches(0.7), row_y, Inches(0.12), Inches(0.75), TEAL)
    add_text(s, Inches(0.95), row_y + Inches(0.18), Inches(4.1), Inches(0.5),
             k, size=16, bold=True, color=NAVY)
    add_text(s, Inches(5.2), row_y + Inches(0.18), Inches(7.8), Inches(0.5),
             v, size=16, color=GREY)

add_rect(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.6), NAVY)
add_text(s, Inches(0.95), Inches(5.55), Inches(11.5), Inches(0.5),
         "Gap addressed", size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.95), Inches(5.9), Inches(11.5), Inches(1.0),
         ("No end-to-end multi-agent pipeline integrates detection, "
          "test generation, and repair with heterogeneous models.\n"
          "ASQA combines LangGraph orchestration + role-specialised LLMs "
          "across the full QA lifecycle."),
         size=16, color=WHITE)

# ---------- 6. Architecture Overview ----------
s = prs.slides.add_slide(BLANK)
header(s, "ASQA Architecture", "Five specialised agents on a LangGraph state machine")

# Pipeline boxes
agents = [
    ("1. Code Reader", "GPT-4.1-mini", "Risk analysis JSON"),
    ("2. Test Generator", "Claude Sonnet 4", "Executable test file"),
    ("3. Runner", "GPT-4.1-mini", "Docker exec + classify"),
    ("4. Bug Reporter", "GPT-4.1-mini", "Structured verdict"),
    ("5. Fix Suggester", "Claude Sonnet 4", "Unified diff patch"),
]
n = len(agents)
gap = Inches(0.15)
total_w = Inches(12.4)
box_w = (total_w - gap * (n - 1)) / n
box_h = Inches(2.2)
y0 = Inches(2.4)
x0 = Inches(0.47)
colors = [TEAL, ACCENT, TEAL, ACCENT, TEAL]
for i, (name, model, out) in enumerate(agents):
    x = x0 + i * (box_w + gap)
    add_rect(s, x, y0, box_w, box_h, WHITE, line=NAVY)
    add_rect(s, x, y0, box_w, Inches(0.35), colors[i])
    add_text(s, x + Inches(0.1), y0 + Inches(0.05), box_w, Inches(0.3),
             name, size=12, bold=True, color=WHITE)
    add_text(s, x + Inches(0.15), y0 + Inches(0.55), box_w - Inches(0.2), Inches(0.4),
             model, size=11, bold=True, color=NAVY)
    add_text(s, x + Inches(0.15), y0 + Inches(0.95), box_w - Inches(0.2), Inches(1.1),
             out, size=11, color=GREY)
    if i < n - 1:
        arrow_x = x + box_w + Emu(int(gap / 4))
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x,
                                 y0 + Inches(0.9), gap, Inches(0.4))
        arr.fill.solid(); arr.fill.fore_color.rgb = NAVY
        arr.line.fill.background()

# Retry loop note
add_text(s, Inches(0.5), Inches(4.9), Inches(12), Inches(0.4),
         "\u21ba  Retry edge: Runner \u2192 Test Generator on execution_error (\u2264 2\u00d7)",
         size=14, bold=True, color=ACCENT)

add_rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.6), LIGHT)
add_text(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.4),
         "Why multi-agent?", size=14, bold=True, color=TEAL)
add_bullets(s, Inches(0.7), Inches(6.0), Inches(12), Inches(1.2), [
    "Isolates each output format \u2192 fewer parse failures",
    "Conditional edges enable reliable retries",
    "Shared TypedDict state keeps context bounded per agent",
], size=14, color=NAVY)

# ---------- 7. Heterogeneous Model Assignment ----------
s = prs.slides.add_slide(BLANK)
header(s, "Heterogeneous Model Assignment",
       "Match the model to the cognitive demand")

cw = Inches(6.0); ch = Inches(4.8); cy = Inches(1.8)
# GPT card
add_rect(s, Inches(0.5), cy, cw, ch, LIGHT)
add_rect(s, Inches(0.5), cy, cw, Inches(0.5), TEAL)
add_text(s, Inches(0.7), cy + Inches(0.08), cw, Inches(0.4),
         "GPT-4.1-mini  (Azure OpenAI)", size=18, bold=True, color=WHITE)
add_text(s, Inches(0.7), cy + Inches(0.7), cw - Inches(0.4), Inches(0.4),
         "Used by:  Code Reader, Runner, Bug Reporter", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), cy + Inches(1.3), cw - Inches(0.4), Inches(3.5), [
    "Fast structured JSON reasoning",
    "Large context for diff processing",
    "Cost-efficient for high-throughput eval",
], size=15)

# Claude card
cx2 = Inches(6.9)
add_rect(s, cx2, cy, cw, ch, LIGHT)
add_rect(s, cx2, cy, cw, Inches(0.5), ACCENT)
add_text(s, cx2 + Inches(0.2), cy + Inches(0.08), cw, Inches(0.4),
         "Claude Sonnet 4  (Anthropic)", size=18, bold=True, color=WHITE)
add_text(s, cx2 + Inches(0.2), cy + Inches(0.7), cw - Inches(0.4), Inches(0.4),
         "Used by:  Test Generator, Fix Suggester", size=14, bold=True, color=NAVY)
add_bullets(s, cx2 + Inches(0.2), cy + Inches(1.3), cw - Inches(0.4), Inches(3.5), [
    "High-quality long-form code generation",
    "Extended thinking for root-cause analysis",
    "Better at minimal, targeted patches",
], size=15)

add_text(s, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
         "Design principle: exploit complementary strengths instead of a single-model compromise.",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ---------- 8. Data & Evaluation Setup ----------
s = prs.slides.add_slide(BLANK)
header(s, "Data & Evaluation Setup")

# Left card
add_rect(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(5.3), LIGHT)
add_text(s, Inches(0.7), Inches(1.75), Inches(6), Inches(0.4),
         "Curated corpus (1,829 bugs)", size=16, bold=True, color=TEAL)
add_bullets(s, Inches(0.7), Inches(2.25), Inches(5.8), Inches(2.8), [
    "BugsInPy \u2014 501 Python bugs",
    "Defects4J \u2014 828 Java bugs",
    "SWE-bench Verified \u2014 500 GitHub issues",
], size=15)
add_text(s, Inches(0.7), Inches(4.3), Inches(6), Inches(0.4),
         "Evaluated end-to-end", size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(0.7), Inches(4.75), Inches(5.8), Inches(2.0), [
    "1,500-bug BugsInPy sample",
    "1,350 Keras  +  150 thefuck",
    "Single pip-installable runtime",
], size=15)

# Right card - metrics
add_rect(s, Inches(6.9), Inches(1.6), Inches(6.0), Inches(5.3), LIGHT)
add_text(s, Inches(7.1), Inches(1.75), Inches(6), Inches(0.4),
         "Evaluation metrics", size=16, bold=True, color=TEAL)
metrics = [
    ("Bug Detection Rate", "% of bugs correctly flagged", "> 60%"),
    ("Mean Time to Report", "pipeline invocation \u2192 output", "< 120 s"),
    ("Completion Rate", "runs with valid output", "\u2191"),
]
for i, (m, d, t) in enumerate(metrics):
    yy = Inches(2.3) + i * Inches(1.4)
    add_rect(s, Inches(7.1), yy, Inches(5.6), Inches(1.2), WHITE)
    add_rect(s, Inches(7.1), yy, Inches(0.12), Inches(1.2), ACCENT)
    add_text(s, Inches(7.35), yy + Inches(0.1), Inches(3.5), Inches(0.4),
             m, size=14, bold=True, color=NAVY)
    add_text(s, Inches(7.35), yy + Inches(0.5), Inches(4), Inches(0.5),
             d, size=12, color=GREY)
    add_text(s, Inches(11.1), yy + Inches(0.35), Inches(1.5), Inches(0.5),
             t, size=18, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)

# ---------- 9. Headline Results ----------
s = prs.slides.add_slide(BLANK)
header(s, "Headline Results", "ASQA on the 1,500-bug sample")
kpis = [
    ("80%", "Bug Detection Rate", "Target > 60%", TEAL),
    ("50.7 s", "Mean Time to Report", "Target < 120 s", ACCENT),
    ("100%", "Completion Rate", "0 failed runs", NAVY),
    ("5", "Artefacts / bug", "vs. 1 for baselines", TEAL),
]
kw = Inches(2.9); kh = Inches(2.6); gx = Inches(0.25)
total_k = kw * 4 + gx * 3
start_x = (SW - total_k) // 2
for i, (v, lbl, sub, c) in enumerate(kpis):
    x = start_x + i * (kw + gx)
    y = Inches(1.9)
    add_rect(s, x, y, kw, kh, WHITE, line=c)
    add_rect(s, x, y, kw, Inches(0.18), c)
    add_text(s, x, y + Inches(0.55), kw, Inches(1.1),
             v, size=54, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.65), kw, Inches(0.4),
             lbl, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(2.05), kw, Inches(0.4),
             sub, size=13, color=GREY, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.7), Inches(4.9), Inches(12), Inches(2.2), LIGHT)
add_text(s, Inches(0.95), Inches(5.0), Inches(12), Inches(0.4),
         "Reading the numbers", size=14, bold=True, color=TEAL)
add_bullets(s, Inches(0.95), Inches(5.45), Inches(12), Inches(1.7), [
    "80% BDR beats the 60% target by +20 pts",
    "Median 50.2 s, 90th percentile 67.7 s, range 30.6\u201386.0 s",
    "Retry loop activated on 100% of runs \u2014 essential reliability mechanism",
], size=15)

# ---------- 10. Outcome Composition (figure) ----------
s = prs.slides.add_slide(BLANK)
header(s, "Outcome Composition vs. Baselines")
if (FIG / "fig_outcomes.png").exists():
    s.shapes.add_picture(str(FIG / "fig_outcomes.png"),
                         Inches(0.4), Inches(1.4), height=Inches(5.6))
add_rect(s, Inches(8.4), Inches(1.8), Inches(4.6), Inches(5.0), LIGHT)
add_text(s, Inches(8.6), Inches(1.95), Inches(4.4), Inches(0.4),
         "What the chart shows", size=14, bold=True, color=TEAL)
add_bullets(s, Inches(8.6), Inches(2.4), Inches(4.3), Inches(4.4), [
    "ASQA: 1,200 detections + 300 calibrated no_bug verdicts",
    "GPT single-agent: flags every input \u2014 no discrimination",
    "Claude single-agent: 750 runs lost to JSON parse failures",
    "Only ASQA produces structured, usable artefacts",
], size=13)

# ---------- 11. Latency distribution (figure) ----------
s = prs.slides.add_slide(BLANK)
header(s, "Mean Time to Report", "All systems below the 120 s target")
if (FIG / "fig_mttr_distribution.png").exists():
    s.shapes.add_picture(str(FIG / "fig_mttr_distribution.png"),
                         Inches(0.4), Inches(1.4), height=Inches(5.6))
add_rect(s, Inches(8.4), Inches(1.8), Inches(4.6), Inches(5.0), LIGHT)
add_text(s, Inches(8.6), Inches(1.95), Inches(4.4), Inches(0.4),
         "Latency trade-off", size=14, bold=True, color=TEAL)
add_bullets(s, Inches(8.6), Inches(2.4), Inches(4.3), Inches(4.4), [
    "5 sequential LLM calls + Docker exec",
    "~3\u00d7 slower than single-agent, still well under budget",
    "Cost of 5 rich artefacts, not 1 opaque blob",
    "Top optimisation targets: Runner + Test Generator",
], size=13)

# ---------- 12. Agent breakdown (figure) ----------
s = prs.slides.add_slide(BLANK)
header(s, "Where Does the Time Go?", "Per-agent contribution to the 50.7 s MTTR")
if (FIG / "fig_agent_breakdown.png").exists():
    s.shapes.add_picture(str(FIG / "fig_agent_breakdown.png"),
                         Inches(0.4), Inches(1.4), height=Inches(5.6))
add_rect(s, Inches(8.4), Inches(1.8), Inches(4.6), Inches(5.0), LIGHT)
add_text(s, Inches(8.6), Inches(1.95), Inches(4.4), Inches(0.4),
         "Bottlenecks", size=14, bold=True, color=TEAL)
add_bullets(s, Inches(8.6), Inches(2.4), Inches(4.3), Inches(4.4), [
    "Runner dominates \u2014 Docker cold-start + retry",
    "Test Generator second \u2014 long-form code output",
    "Runner + TG \u2248 55% of wall-clock time",
    "Parallel Code Reader \u2015 TG could shave seconds",
], size=13)

# ---------- 13. Statistical Significance ----------
s = prs.slides.add_slide(BLANK)
header(s, "Statistical Significance",
       "McNemar's paired exact test on per-bug verdicts")

# Table
tbl_x = Inches(0.7); tbl_y = Inches(1.7)
cols = [Inches(4.8), Inches(1.7), Inches(1.7), Inches(1.9)]
hdr = ["Comparison", "b\u2081\u2080", "b\u2080\u2081", "Exact p"]
add_rect(s, tbl_x, tbl_y, sum(cols, Emu(0)), Inches(0.55), NAVY)
x = tbl_x
for i, h in enumerate(hdr):
    add_text(s, x + Inches(0.15), tbl_y + Inches(0.12), cols[i], Inches(0.4),
             h, size=16, bold=True, color=WHITE)
    x += cols[i]
rows = [
    ("ASQA vs. Claude SL", "600", "150", "< 10\u207b\u2076"),
    ("ASQA vs. GPT SL",    "0",   "300", "< 10\u207b\u2076"),
]
for r_idx, row in enumerate(rows):
    ry = tbl_y + Inches(0.55) + r_idx * Inches(0.55)
    fill = LIGHT if r_idx % 2 == 0 else WHITE
    add_rect(s, tbl_x, ry, sum(cols, Emu(0)), Inches(0.55), fill)
    x = tbl_x
    for i, v in enumerate(row):
        add_text(s, x + Inches(0.15), ry + Inches(0.13), cols[i], Inches(0.4),
                 v, size=15, color=NAVY, bold=(i == 0))
        x += cols[i]

add_rect(s, Inches(0.7), Inches(4.2), Inches(12), Inches(2.8), LIGHT)
add_text(s, Inches(0.9), Inches(4.35), Inches(12), Inches(0.4),
         "Interpretation", size=14, bold=True, color=TEAL)
add_bullets(s, Inches(0.9), Inches(4.8), Inches(12), Inches(2.2), [
    "ASQA's 100% vs. 50% completion rate over Claude-SL is not chance.",
    "GPT-SL's higher BDR is also significant \u2014 but indiscriminately flags "
    "everything (zero no_bug verdicts).",
    "Latency gap is uniformly positive \u2192 sign-test p < 10\u207b\u2076.",
], size=14)

# ---------- 14. Scalability (figure) ----------
s = prs.slides.add_slide(BLANK)
header(s, "Scalability & Throughput")
if (FIG / "fig_throughput.png").exists():
    s.shapes.add_picture(str(FIG / "fig_throughput.png"),
                         Inches(0.4), Inches(1.4), height=Inches(5.6))
add_rect(s, Inches(8.4), Inches(1.8), Inches(4.6), Inches(5.0), LIGHT)
add_text(s, Inches(8.6), Inches(1.95), Inches(4.4), Inches(0.4),
         "Throughput", size=14, bold=True, color=TEAL)
add_bullets(s, Inches(8.6), Inches(2.4), Inches(4.3), Inches(4.4), [
    "ASQA: 71 bugs/hour (single worker)",
    "GPT-SL: 219 / hour  \u2022  Claude-SL: 131 / hour",
    "Linear \u2014 constant per-bug cost",
    "Read-only state \u2192 trivial horizontal scaling",
    "~3\u20134 h on 8 workers for full 1,829 corpus",
], size=13)

# ---------- 15. Error Analysis ----------
s = prs.slides.add_slide(BLANK)
header(s, "Error Analysis", "What the numbers actually mean")

cards = [
    ("ASQA \u2014 20% no_bug",
     "Reflects calibrated confidence. Many non-detections are "
     "stylistic diffs or deprecated APIs, not functional faults.",
     TEAL),
    ("Claude SL \u2014 50% failures",
     "Mega-prompt fragility: unclosed brackets, trailing commas, "
     "mixed quote styles break JSON parsing.",
     ACCENT),
    ("GPT SL \u2014 100% BDR",
     "Flags every input as a bug. High recall, low precision \u2014 "
     "cannot discriminate benign from faulty diffs.",
     NAVY),
    ("Retry loop \u2014 essential",
     "Every run needed exactly one retry: Docker sandbox first-attempt "
     "errors reliably recovered by adjusted tests.",
     TEAL),
]
for i, (t, d, c) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.2)
    y = Inches(1.6) + row * Inches(2.7)
    add_rect(s, x, y, Inches(6.0), Inches(2.5), LIGHT)
    add_rect(s, x, y, Inches(0.15), Inches(2.5), c)
    add_text(s, x + Inches(0.35), y + Inches(0.2), Inches(5.5), Inches(0.5),
             t, size=16, bold=True, color=NAVY)
    add_text(s, x + Inches(0.35), y + Inches(0.8), Inches(5.5), Inches(1.6),
             d, size=14, color=GREY)

# ---------- 16. Implications ----------
s = prs.slides.add_slide(BLANK)
header(s, "Project Implications",
       "What this tells us about autonomous QA design")
items = [
    ("Specialisation beats monolith",
     "Five focused prompts yield richer, parseable outputs than one mega-prompt."),
    ("Heterogeneous models > single model",
     "GPT for structured reasoning + Claude for code gen avoids per-task compromise."),
    ("Orchestration is a reliability tool",
     "LangGraph conditional edges turn transient errors into recoveries, not failures."),
    ("Multi-stage evidence = calibration",
     "Cumulative artefacts let the system abstain, unlike single-prompt baselines."),
    ("Practical latency budget",
     "50.7 s MTTR is compatible with pre-merge CI / PR review workflows."),
]
for i, (h, d) in enumerate(items):
    y = Inches(1.6) + i * Inches(1.05)
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.95), LIGHT if i % 2 == 0 else WHITE)
    add_rect(s, Inches(0.6), y, Inches(0.15), Inches(0.95), TEAL)
    add_text(s, Inches(0.9), y + Inches(0.12), Inches(4.4), Inches(0.5),
             h, size=16, bold=True, color=NAVY)
    add_text(s, Inches(5.5), y + Inches(0.22), Inches(7.2), Inches(0.6),
             d, size=14, color=GREY)

# ---------- 17. Limitations ----------
s = prs.slides.add_slide(BLANK)
header(s, "Limitations", "Honest caveats on the current study")
lims = [
    ("Single-language evaluation",
     "End-to-end run on Python / BugsInPy only. Defects4J (Java/Maven) and "
     "SWE-bench Verified curated but not executed \u2014 build-env engineering gap."),
    ("Project concentration",
     "90% of the sample is Keras \u2014 detection may be inflated by LLM prior exposure."),
    ("Throughput ceiling",
     "Docker cold-start + local checkout cap single-worker throughput at 71 bugs/h."),
    ("Model availability",
     "Azure deployment constraints forced GPT-4.1-mini over GPT-4o for GPT agents."),
    ("FPR not fully automated",
     "False-positive assessment needs human annotators with Cohen's Kappa \u2014 "
     "not done in this study."),
]
for i, (h, d) in enumerate(lims):
    y = Inches(1.5) + i * Inches(1.05)
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.95), LIGHT if i % 2 == 0 else WHITE)
    add_rect(s, Inches(0.6), y, Inches(0.15), Inches(0.95), ACCENT)
    add_text(s, Inches(0.9), y + Inches(0.12), Inches(4.2), Inches(0.5),
             h, size=15, bold=True, color=NAVY)
    add_text(s, Inches(5.3), y + Inches(0.1), Inches(7.4), Inches(0.9),
             d, size=13, color=GREY)

# ---------- 18. Future Work ----------
s = prs.slides.add_slide(BLANK)
header(s, "Future Work")
future = [
    ("Cross-language evaluation",
     "Extend to Defects4J & SWE-bench for Java + real-world issue resolution."),
    ("Model upgrades",
     "GPT-4o for reasoning agents; measure BDR and precision lift."),
    ("Parallel agents",
     "Run Code Reader and Test Generator concurrently where deps allow."),
    ("CI/CD integration",
     "Wire into pull-request review \u2014 real-time bug / fix suggestions."),
    ("Human evaluation",
     "False-Positive Rate with inter-annotator agreement (Cohen's Kappa)."),
]
for i, (h, d) in enumerate(future):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.2)
    y = Inches(1.6) + row * Inches(1.7)
    add_rect(s, x, y, Inches(6.0), Inches(1.5), LIGHT)
    add_rect(s, x, y, Inches(0.15), Inches(1.5), TEAL)
    add_text(s, x + Inches(0.35), y + Inches(0.15), Inches(5.5), Inches(0.5),
             h, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.35), y + Inches(0.65), Inches(5.5), Inches(0.8),
             d, size=13, color=GREY)

# ---------- 19. Conclusion ----------
s = prs.slides.add_slide(BLANK)
header(s, "Conclusion")
concl = [
    "Multi-agent decomposition works: 80% BDR, 100% completion.",
    "Heterogeneous model assignment outperforms any single model.",
    "LangGraph conditional edges turn retries into a built-in reliability layer.",
    "Latency overhead (~3\u00d7) buys five rich artefacts per bug \u2014 worth it.",
    "Scope caveat: Python-only evaluation; Java & SWE-bench are next.",
]
add_rect(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(4.6), LIGHT)
for i, line in enumerate(concl):
    y = Inches(1.95) + i * Inches(0.9)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.1),
                              Inches(0.35), Inches(0.35))
    circ.fill.solid(); circ.fill.fore_color.rgb = TEAL
    circ.line.fill.background()
    add_text(s, Inches(1.6), y + Inches(0.05), Inches(10.8), Inches(0.7),
             line, size=17, color=NAVY)

add_rect(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.7), NAVY)
add_text(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.4),
         "ASQA delivers autonomous QA that is measurable, reliable, and extensible.",
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ---------- 20. Thank You / Q&A ----------
s = prs.slides.add_slide(BLANK)
add_rect(s, Emu(0), Emu(0), SW, SH, NAVY)
add_rect(s, Emu(0), Inches(3.3), SW, Inches(0.06), TEAL)
add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(1.2),
         "Thank you", size=72, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.6), Inches(12), Inches(0.6),
         "Questions?", size=32, color=RGBColor(0xCF, 0xE3, 0xE7))
add_text(s, Inches(0.7), Inches(5.3), Inches(12), Inches(0.4),
         "Smruti Pote   \u2022   x24269522@student.ncirl.ie", size=18, color=WHITE)
add_text(s, Inches(0.7), Inches(5.75), Inches(12), Inches(0.4),
         "MSc Artificial Intelligence  \u2014  National College of Ireland",
         size=14, color=RGBColor(0xCF, 0xE3, 0xE7))

prs.save(str(OUT))
print(f"Saved: {OUT}")
